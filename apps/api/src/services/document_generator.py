"""Generate Office documents and PDFs as in-memory bytes."""

from __future__ import annotations

from io import BytesIO
from typing import Any

from docx import Document
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Cm, Pt, RGBColor as DocxRGBColor
from openpyxl import Workbook
from openpyxl.styles import Alignment, Font, PatternFill
from openpyxl.utils import get_column_letter
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt as PptPt

NAVY = RGBColor(15, 23, 42)
GOLD = RGBColor(212, 160, 23)
WHITE = RGBColor(255, 255, 255)
SLATE = RGBColor(30, 41, 59)


def create_document(
    doc_type: str,
    template: str | None,
    title: str,
    context: dict[str, Any] | None = None,
) -> bytes:
    """Create a DOCX, XLSX, PPTX, or PDF document as bytes."""
    normalized_template = template or "blank"
    ctx = _build_context(title, context)

    if doc_type == "docx":
        return _create_docx(normalized_template, title, ctx)
    if doc_type == "xlsx":
        return _create_xlsx(normalized_template, title, ctx)
    if doc_type == "pptx":
        return _create_pptx(normalized_template, title, ctx)
    if doc_type == "pdf":
        return _create_pdf(normalized_template, title, ctx)
    raise ValueError(f"Unsupported office document type: {doc_type}")


def _build_context(title: str, context: dict[str, Any] | None) -> dict[str, Any]:
    ctx = dict(context or {})
    ctx.setdefault("title", title)
    ctx.setdefault("matter", title)
    ctx.setdefault("recipient", "[Recipient]")
    ctx.setdefault("author", "Amin Legal Operations")
    ctx.setdefault("date", "[Date]")
    ctx.setdefault("client_name", "[Client]")
    ctx.setdefault("counterparty", "[Counterparty]")
    ctx.setdefault("court_name", "[Court / Forum]")
    ctx.setdefault("city", "Saudi Arabia")
    ctx.setdefault(
        "sections",
        [
            "Presentation draft generated for demo purposes only.",
            "Replace with live matter details before external circulation.",
        ],
    )
    return ctx


def _create_docx(template: str, title: str, ctx: dict[str, Any]) -> bytes:
    doc = Document()
    _apply_docx_defaults(doc)

    doc.add_heading(title, level=0)
    if ctx.get("client_name"):
        doc.add_paragraph(f"Client: {ctx['client_name']}")
    if ctx.get("client_name_ar"):
        doc.add_paragraph(f"Arabic name: {ctx['client_name_ar']}")
    doc.add_paragraph()

    if template == "blank":
        for paragraph in ctx.get("sections", []):
            doc.add_paragraph(str(paragraph))
    elif template == "legal_memo":
        _add_labeled_paragraph(doc, "To", str(ctx.get("recipient")))
        _add_labeled_paragraph(doc, "From", str(ctx.get("author")))
        _add_labeled_paragraph(doc, "Date", str(ctx.get("date")))
        _add_labeled_paragraph(doc, "Re", str(ctx.get("matter")))
        _add_section_block(
            doc,
            "Executive Summary",
            [
                f"This memorandum relates to {ctx.get('matter')} for {ctx.get('client_name')}.",
                "It is generated for presentation use and can be refined in the editor.",
            ],
        )
        _add_section_block(doc, "Facts", _as_string_list(ctx.get("sections")))
        _add_section_block(
            doc,
            "Analysis",
            [
                f"Saudi law and procedural considerations should be aligned to {ctx.get('city')}.",
                f"Counterparty / stakeholder reference: {ctx.get('counterparty')}.",
            ],
        )
        _add_section_block(
            doc,
            "Recommendation",
            [
                "Prioritize evidence discipline, bilingual alignment, and executive-ready next steps.",
            ],
        )
    elif template == "contract":
        _add_section_block(
            doc,
            "Parties",
            [
                f"{ctx.get('client_name')} and {ctx.get('counterparty')}.",
                "This draft is prepared for demonstration purposes only.",
            ],
        )
        _add_section_block(doc, "Commercial Background", _as_string_list(ctx.get("sections")))
        for heading in [
            "Definitions",
            "Scope and Deliverables",
            "Fees and Payment",
            "Confidentiality",
            "Representations and Warranties",
            "Indemnities",
            "Term and Termination",
            "Governing Law and Dispute Resolution",
        ]:
            _add_section_block(doc, heading, [f"Drafting placeholder for {heading.lower()}."])
    elif template == "court_brief":
        _add_labeled_paragraph(doc, "Matter", str(ctx.get("matter")))
        _add_labeled_paragraph(doc, "Claimant", str(ctx.get("client_name")))
        _add_labeled_paragraph(doc, "Respondent", str(ctx.get("counterparty")))
        _add_labeled_paragraph(doc, "Court", str(ctx.get("court_name")))
        _add_section_block(doc, "Relief Sought", ["Orders requested from the competent court or forum."])
        _add_section_block(doc, "Statement of Facts", _as_string_list(ctx.get("sections")))
        _add_section_block(
            doc,
            "Legal Grounds",
            [
                f"Forum strategy should remain consistent with {ctx.get('city')} procedural practice.",
                "Support each point with clean bilingual references where appropriate.",
            ],
        )
        _add_section_block(doc, "Attachments", ["Exhibits, correspondence, and translated schedules."])
    elif template == "demand_letter":
        _add_labeled_paragraph(doc, "To", str(ctx.get("counterparty")))
        _add_labeled_paragraph(doc, "On behalf of", str(ctx.get("client_name")))
        _add_section_block(doc, "Without Prejudice Background", _as_string_list(ctx.get("sections")))
        _add_section_block(doc, "Demand", ["Please remedy the outstanding default within the stated deadline."])
        _add_section_block(doc, "Reservation of Rights", ["All rights and remedies are expressly reserved."])
    elif template == "engagement_letter":
        _add_labeled_paragraph(doc, "Client", str(ctx.get("client_name")))
        _add_labeled_paragraph(doc, "Matter", str(ctx.get("matter")))
        _add_section_block(doc, "Scope of Engagement", _as_string_list(ctx.get("sections")))
        _add_section_block(doc, "Team and Reporting", ["Lead counsel, escalation path, and bilingual reporting expectations."])
        _add_section_block(doc, "Fees and Assumptions", ["Billing assumptions, disbursements, and approval thresholds."])
    else:
        _add_section_block(doc, "Summary", _as_string_list(ctx.get("sections")))

    buffer = BytesIO()
    doc.save(buffer)
    return buffer.getvalue()


def _apply_docx_defaults(doc: Document) -> None:
    for style_name in ("Normal", "Body Text", "Default Paragraph Font"):
        try:
            style = doc.styles[style_name]
        except KeyError:
            continue
        style.font.name = "Calibri"
        style.font.size = Pt(11)
        if style._element.rPr is None:
            style._element.get_or_add_rPr()
        r_fonts = style._element.rPr.rFonts
        if r_fonts is None:
            r_fonts = OxmlElement("w:rFonts")
            style._element.rPr.append(r_fonts)
        r_fonts.set(qn("w:ascii"), "Calibri")
        r_fonts.set(qn("w:hAnsi"), "Calibri")
        r_fonts.set(qn("w:eastAsia"), "Calibri")

    for section in doc.sections:
        section.top_margin = Cm(2.5)
        section.bottom_margin = Cm(2.5)
        section.left_margin = Cm(2.5)
        section.right_margin = Cm(2.5)

    _add_demo_header(doc)
    doc.core_properties.language = "en-US"


def _add_demo_header(doc: Document) -> None:
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    for section in doc.sections:
        header = section.header
        header.is_linked_to_previous = False
        para = header.paragraphs[0] if header.paragraphs else header.add_paragraph()
        para.text = ""
        run = para.add_run("Created with HeyAmin - Demo")
        run.font.size = Pt(8)
        run.font.color.rgb = DocxRGBColor(180, 180, 180)
        run.font.name = "Calibri"
        para.alignment = WD_ALIGN_PARAGRAPH.RIGHT


def _add_labeled_paragraph(doc: Document, label: str, value: str) -> None:
    paragraph = doc.add_paragraph()
    label_run = paragraph.add_run(f"{label}: ")
    label_run.bold = True
    paragraph.add_run(value)


def _add_section_block(doc: Document, heading: str, paragraphs: list[str]) -> None:
    doc.add_heading(heading, level=1)
    for paragraph in paragraphs:
        doc.add_paragraph(paragraph)


def _create_xlsx(template: str, title: str, ctx: dict[str, Any]) -> bytes:
    wb = Workbook()
    ws = wb.active
    ws.title = "Overview"

    headers = [str(item) for item in ctx.get("headers", [])]
    rows = [
        ["" if cell is None else cell for cell in row]
        for row in ctx.get("rows", [])
        if isinstance(row, list)
    ]

    if template == "blank":
        ws["A1"] = title
        ws["A2"] = "Start working in this sheet."
    elif template == "tracker":
        _write_sheet(
            ws,
            headers or ["Title", "Status", "Owner", "Due Date", "Notes"],
            rows
            or [
                ["Prepare filing", "In Progress", "Case Team", "[YYYY-MM-DD]", ""],
                ["Collect exhibits", "Open", "Client", "[YYYY-MM-DD]", ""],
                ["Client approval", "Pending", "Partner", "[YYYY-MM-DD]", ""],
            ],
        )
    elif template == "legal_matrix":
        _write_sheet(
            ws,
            headers or ["Jurisdiction", "Law", "Article", "Summary", "Status"],
            rows
            or [
                ["KSA", "[Law Name]", "[Article]", "[Summary]", "Action"],
                ["KSA", "[Law Name]", "[Article]", "[Summary]", "Monitor"],
            ],
        )
    elif template == "chronology":
        _write_sheet(
            ws,
            headers or ["Date", "Action", "Owner", "Status", "Comment"],
            rows
            or [
                ["2026-04-01", "Matter opened", "Case Team", "Done", ""],
                ["2026-04-05", "Documents requested", "Client", "Done", ""],
                ["2026-04-10", "Next submission", "External Counsel", "Open", ""],
            ],
        )
    elif template == "fee_schedule":
        _write_sheet(
            ws,
            headers or ["Line Item", "Claim", "Response", "Delta", "Comment"],
            rows
            or [
                ["Fees", "0", "0", "0", ""],
                ["Disbursements", "0", "0", "0", ""],
            ],
        )
    elif template == "evidence_index":
        _write_sheet(
            ws,
            headers or ["Exhibit", "Description", "Owner", "Status", "Notes"],
            rows
            or [
                ["C-01", "Contract", "Case Team", "Ready", ""],
                ["C-02", "Correspondence", "Client", "Pending", ""],
            ],
        )
    else:
        _write_sheet(ws, headers or ["Section", "Detail"], rows or [["Summary", title]])

    _style_workbook(wb)
    buffer = BytesIO()
    wb.save(buffer)
    return buffer.getvalue()


def _write_sheet(ws, headers: list[str], rows: list[list[object]]) -> None:
    ws.append(headers)
    for row in rows:
        ws.append(row)
    ws.freeze_panes = "A2"
    ws.auto_filter.ref = ws.dimensions


def _style_workbook(wb: Workbook) -> None:
    header_fill = PatternFill(fill_type="solid", fgColor="1E293B")
    header_font = Font(color="FFFFFF", bold=True)
    bold_font = Font(bold=True)

    for ws in wb.worksheets:
        for cell in ws[1]:
            cell.fill = header_fill
            cell.font = header_font
            cell.alignment = Alignment(horizontal="center", vertical="center")

        for row in ws.iter_rows(min_row=1, max_row=ws.max_row, min_col=1, max_col=ws.max_column):
            for cell in row:
                if cell.row > 1 and cell.column == 1:
                    cell.font = bold_font
                cell.alignment = Alignment(vertical="top", wrap_text=True)

        for idx in range(1, ws.max_column + 1):
            max_length = 0
            for cell in ws[get_column_letter(idx)]:
                max_length = max(max_length, len(str(cell.value or "")))
            ws.column_dimensions[get_column_letter(idx)].width = min(
                max(max_length + 2, 12),
                36,
            )


def _create_pptx(template: str, title: str, ctx: dict[str, Any]) -> bytes:
    prs = Presentation()

    if template == "blank":
        _add_title_slide(prs, title, "Prepared with Amin")
    elif template == "pitch":
        _add_title_slide(prs, title, "Executive pitch")
        for heading in ["Problem", "Solution", "Market", "Team", "Ask"]:
            _add_content_slide(prs, heading, ["[Insert content]", "[Key takeaway]"])
    elif template == "status_update":
        _add_title_slide(prs, title, "Status update")
        slides = ctx.get("slides")
        if isinstance(slides, list) and slides:
            for slide in slides:
                if isinstance(slide, tuple) and len(slide) == 2:
                    _add_content_slide(prs, str(slide[0]), _as_string_list(slide[1]))
        else:
            for heading in ["Summary", "Progress", "Risks", "Next Steps"]:
                _add_content_slide(prs, heading, ["[Insert update]", "[Owner / timeline]"])
    elif template == "legal_overview":
        _add_title_slide(prs, title, "Legal overview")
        for heading in ["Jurisdiction", "Key Laws", "Obligations", "Recommendations"]:
            _add_content_slide(prs, heading, ["[Insert legal point]", "[Commercial implication]"])
    elif template == "board_pack":
        _add_title_slide(prs, title, "Board pack")
        slides = ctx.get("slides")
        if isinstance(slides, list) and slides:
            for slide in slides:
                if isinstance(slide, tuple) and len(slide) == 2:
                    _add_content_slide(prs, str(slide[0]), _as_string_list(slide[1]))
        else:
            for heading in ["Transaction Snapshot", "Key Risks", "Approvals", "Decision Request"]:
                _add_content_slide(prs, heading, ["[Insert board-ready point]"])
    else:
        _add_title_slide(prs, title, "Blank presentation")

    buffer = BytesIO()
    prs.save(buffer)
    return buffer.getvalue()


def _add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = NAVY

    slide.shapes.title.text = title
    title_frame = slide.shapes.title.text_frame
    title_frame.paragraphs[0].font.color.rgb = WHITE
    title_frame.paragraphs[0].font.size = PptPt(28)
    title_frame.paragraphs[0].font.bold = True

    subtitle_box = slide.placeholders[1]
    subtitle_box.text = subtitle
    subtitle_para = subtitle_box.text_frame.paragraphs[0]
    subtitle_para.font.color.rgb = GOLD
    subtitle_para.font.size = PptPt(16)
    subtitle_para.alignment = PP_ALIGN.LEFT


def _add_content_slide(prs: Presentation, title: str, bullets: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    slide.shapes.title.text = title
    title_para = slide.shapes.title.text_frame.paragraphs[0]
    title_para.font.color.rgb = NAVY
    title_para.font.bold = True
    title_para.font.size = PptPt(22)

    body = slide.placeholders[1].text_frame
    body.clear()
    for idx, bullet in enumerate(bullets):
        paragraph = body.paragraphs[0] if idx == 0 else body.add_paragraph()
        paragraph.text = bullet
        paragraph.font.color.rgb = SLATE
        paragraph.font.size = PptPt(16)
        paragraph.level = 0

    accent = slide.shapes.add_shape(
        1,
        Inches(0.5),
        Inches(1.3),
        Inches(1.5),
        Inches(0.08),
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = GOLD
    accent.line.fill.background()


def _create_pdf(template: str, title: str, ctx: dict[str, Any]) -> bytes:
    lines = _pdf_lines(template, title, ctx)
    try:
        import fitz

        doc = fitz.open()
        page = doc.new_page(width=595, height=842)
        rect = fitz.Rect(54, 72, 540, 790)
        page.insert_textbox(
            rect,
            "\n".join(lines),
            fontsize=10.5,
            fontname="helv",
            color=(0.12, 0.16, 0.23),
            align=fitz.TEXT_ALIGN_LEFT,
        )

        pdf_bytes = doc.tobytes(garbage=4, deflate=True)
        doc.close()
        return pdf_bytes
    except ImportError:
        return _create_fallback_pdf(lines)


def _pdf_lines(template: str, title: str, ctx: dict[str, Any]) -> list[str]:
    header = [
        title,
        "Prepared for presentation purposes only",
        f"Client: {ctx.get('client_name')}",
        f"Matter: {ctx.get('matter')}",
    ]

    if template == "filing_receipt":
        body = [
            f"Forum / court: {ctx.get('court_name')}",
            f"Counterparty: {ctx.get('counterparty')}",
            f"City: {ctx.get('city')}",
            "",
            "Receipt summary:",
            "1. Electronic filing reference generated for demo display.",
            "2. Supporting bundle logged against the matter workspace.",
            "3. Arabic and English references aligned for presentation use.",
        ]
    elif template == "client_notice":
        body = [
            f"Recipient summary prepared for {ctx.get('client_name')}.",
            f"Primary reference: {ctx.get('counterparty')}",
            "",
            "Key points:",
        ] + [f"- {item}" for item in _as_string_list(ctx.get("sections"))]
    elif template == "hearing_minutes":
        body = [
            f"Hearing / meeting note for {ctx.get('matter')}.",
            "",
            "Minutes:",
        ] + [f"- {item}" for item in _as_string_list(ctx.get("sections"))]
    elif template == "settlement_terms":
        body = [
            f"Settlement summary for {ctx.get('matter')}.",
            "",
            "Commercial terms:",
        ] + [f"- {item}" for item in _as_string_list(ctx.get("sections"))]
    else:
        body = [f"- {item}" for item in _as_string_list(ctx.get("sections"))]

    return header + [""] + body


def _as_string_list(value: Any) -> list[str]:
    if isinstance(value, list):
        return [str(item) for item in value]
    if value is None:
        return []
    return [str(value)]


def _create_fallback_pdf(lines: list[str]) -> bytes:
    escaped_lines = [
        line.replace("\\", "\\\\").replace("(", "\\(").replace(")", "\\)")
        for line in lines
    ]
    commands = ["BT", "/F1 10 Tf", "54 780 Td", "14 TL"]
    for index, line in enumerate(escaped_lines):
        if index == 0:
            commands.append(f"({line}) Tj")
        else:
            commands.append(f"T* ({line}) Tj")
    commands.append("ET")
    stream = "\n".join(commands).encode("latin-1", errors="ignore")

    objects = [
        b"1 0 obj << /Type /Catalog /Pages 2 0 R >> endobj",
        b"2 0 obj << /Type /Pages /Kids [3 0 R] /Count 1 >> endobj",
        (
            b"3 0 obj << /Type /Page /Parent 2 0 R /MediaBox [0 0 595 842] "
            b"/Resources << /Font << /F1 4 0 R >> >> /Contents 5 0 R >> endobj"
        ),
        b"4 0 obj << /Type /Font /Subtype /Type1 /BaseFont /Helvetica >> endobj",
        f"5 0 obj << /Length {len(stream)} >> stream\n".encode("latin-1")
        + stream
        + b"\nendstream endobj",
    ]

    pdf = bytearray(b"%PDF-1.4\n")
    offsets = [0]
    for obj in objects:
        offsets.append(len(pdf))
        pdf.extend(obj)
        pdf.extend(b"\n")

    xref_start = len(pdf)
    pdf.extend(f"xref\n0 {len(offsets)}\n".encode("latin-1"))
    pdf.extend(b"0000000000 65535 f \n")
    for offset in offsets[1:]:
        pdf.extend(f"{offset:010d} 00000 n \n".encode("latin-1"))
    pdf.extend(
        (
            f"trailer << /Size {len(offsets)} /Root 1 0 R >>\n"
            f"startxref\n{xref_start}\n%%EOF"
        ).encode("latin-1")
    )
    return bytes(pdf)
