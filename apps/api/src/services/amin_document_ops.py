"""Amin's Office document operations engine."""

from __future__ import annotations

import json
import re
from io import BytesIO
from typing import Any

from docx import Document
from docx.enum.text import WD_BREAK
from openpyxl import load_workbook
from openpyxl.styles import Font, PatternFill
from openpyxl.utils import get_column_letter
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.util import Inches

from src.services.agent.llm_router import chat_completion


async def execute_instruction(
    instruction: str,
    file_bytes: bytes,
    doc_type: str,
    context: dict[str, Any] | None = None,
) -> tuple[bytes, list[dict[str, Any]], str]:
    """Plan and apply Office operations from a natural-language instruction."""
    if doc_type == "pdf":
        raise ValueError("PDF documents are view-only and cannot be edited by Amin.")

    document_state = extract_document_state(file_bytes, doc_type)
    ops = await _plan_operations(
        instruction=instruction,
        document_state=document_state,
        doc_type=doc_type,
        context=context or {},
    )
    new_bytes = _apply_operations(file_bytes=file_bytes, doc_type=doc_type, ops=ops)
    summary = await _summarize_changes(instruction, doc_type, ops)
    return new_bytes, ops, summary


def extract_document_state(file_bytes: bytes, doc_type: str) -> dict[str, Any]:
    """Extract a readable structured representation of the current document."""
    if doc_type == "docx":
        return _extract_docx_state(file_bytes)
    if doc_type == "xlsx":
        return _extract_xlsx_state(file_bytes)
    if doc_type == "pptx":
        return _extract_pptx_state(file_bytes)
    if doc_type == "pdf":
        return _extract_pdf_state(file_bytes)
    raise ValueError(f"Unsupported office document type: {doc_type}")


def extract_document_metadata(file_bytes: bytes, doc_type: str) -> dict[str, Any]:
    """Build preview-friendly metadata from a document."""
    state = extract_document_state(file_bytes, doc_type)
    if doc_type == "docx":
        paragraphs = state.get("paragraphs", [])
        headings = state.get("headings", [])
        word_count = sum(len((item.get("text") or "").split()) for item in paragraphs)
        page_breaks = sum(
            1 for item in paragraphs if "[PAGE_BREAK]" in (item.get("text") or "")
        )
        return {
            "paragraph_count": len(paragraphs),
            "word_count": word_count,
            "page_count": max(page_breaks + 1, 1),
            "headings": headings[:20],
            "preview_paragraphs": [
                item.get("text", "") for item in paragraphs if item.get("text")
            ][:12],
            "preview_text": "\n\n".join(
                item.get("text", "") for item in paragraphs if item.get("text")
            )[:1800],
        }
    if doc_type == "xlsx":
        sheets = state.get("sheets", [])
        first_sheet = sheets[0] if sheets else {}
        return {
            "sheet_count": len(sheets),
            "sheet_names": [sheet.get("name") for sheet in sheets],
            "preview_sheet": first_sheet,
        }
    if doc_type == "pptx":
        slides = state.get("slides", [])
        return {
            "slide_count": len(slides),
            "slide_titles": [
                slide.get("title") or f"Slide {idx + 1}"
                for idx, slide in enumerate(slides)
            ],
            "preview_slides": slides[:8],
        }

    pages = state.get("pages", [])
    preview_text = "\n\n".join(
        page.get("text", "") for page in pages if page.get("text")
    )[:1800]
    return {
        "page_count": state.get("page_count", len(pages)),
        "preview_pages": pages[:6],
        "preview_text": preview_text,
    }


async def answer_document_question(
    file_bytes: bytes,
    doc_type: str,
    question: str | None = None,
) -> str:
    """Answer a question about a document or summarize it."""
    state = extract_document_state(file_bytes, doc_type)
    prompt = json.dumps(state, ensure_ascii=True)
    if not question:
        question = "Summarize this document for the user."
    response = await chat_completion(
        messages=[
            {
                "role": "system",
                "content": (
                    "You are Amin, a legal AI assistant. Answer strictly from the provided "
                    "document state. If the state is sparse, say so clearly."
                ),
            },
            {
                "role": "user",
                "content": f"Question: {question}\n\nDocument state JSON:\n{prompt}",
            },
        ],
        tools=None,
        model="gpt-4o",
    )
    content = getattr(response.choices[0].message, "content", "") or ""
    if content and "stub mode" not in content.lower():
        return content
    return _fallback_document_answer(state, doc_type, question)


async def _plan_operations(
    instruction: str,
    document_state: dict[str, Any],
    doc_type: str,
    context: dict[str, Any],
) -> list[dict[str, Any]]:
    available_ops = {
        "docx": [
            "add_heading(text, level)",
            "add_paragraph(text, style?)",
            "add_table(rows, cols, data)",
            "replace_text(find, replace)",
            "set_style(paragraph_index, style)",
            "add_page_break()",
            "clear_after(paragraph_index)",
            "bold_text(paragraph_index)",
            "set_font(paragraph_index, size, name)",
        ],
        "xlsx": [
            "set_cell(sheet, row, col, value)",
            "set_formula(sheet, row, col, formula)",
            "add_row(sheet, row_index, data)",
            "bold_row(sheet, row)",
            "set_fill(sheet, row, col, hex_color)",
            "add_sheet(name)",
            "rename_sheet(old, new)",
            "auto_fit_columns(sheet)",
        ],
        "pptx": [
            "add_slide(layout_index, title, body?)",
            "edit_slide_title(slide_index, text)",
            "edit_slide_body(slide_index, text)",
            "add_bullet(slide_index, text, level?)",
            "add_image_placeholder(slide_index)",
            "reorder_slides(new_order)",
            "delete_slide(slide_index)",
            "set_slide_notes(slide_index, text)",
        ],
    }[doc_type]

    response = await chat_completion(
        messages=[
            {
                "role": "system",
                "content": (
                    "You are Amin, a legal AI assistant. You are given a document state and "
                    "an instruction. Return a JSON array of operations only. Each operation "
                    "must be an object with a type field and the parameters needed to execute it. "
                    f"Allowed operations for {doc_type}: {', '.join(available_ops)}."
                ),
            },
            {
                "role": "user",
                "content": json.dumps(
                    {
                        "instruction": instruction,
                        "context": context,
                        "document_state": document_state,
                    },
                    ensure_ascii=True,
                ),
            },
        ],
        tools=None,
        model="gpt-4o",
    )
    content = getattr(response.choices[0].message, "content", "") or ""
    parsed = _extract_json_array(content)
    if parsed:
        return parsed
    return _fallback_operations(doc_type, instruction)


def _extract_json_array(text: str) -> list[dict[str, Any]]:
    match = re.search(r"\[[\s\S]*\]", text)
    if not match:
        return []
    try:
        data = json.loads(match.group(0))
    except json.JSONDecodeError:
        return []
    if isinstance(data, list):
        return [item for item in data if isinstance(item, dict)]
    return []


def _fallback_operations(doc_type: str, instruction: str) -> list[dict[str, Any]]:
    if doc_type == "docx":
        return [{"type": "add_paragraph", "text": instruction}]
    if doc_type == "xlsx":
        return [
            {"type": "set_cell", "sheet": "Overview", "row": 2, "col": 1, "value": instruction},
            {"type": "auto_fit_columns", "sheet": "Overview"},
        ]
    return [{"type": "add_slide", "layout_index": 1, "title": "Amin Update", "body": instruction}]


def _apply_operations(file_bytes: bytes, doc_type: str, ops: list[dict[str, Any]]) -> bytes:
    if doc_type == "docx":
        return _apply_docx_operations(file_bytes, ops)
    if doc_type == "xlsx":
        return _apply_xlsx_operations(file_bytes, ops)
    if doc_type == "pptx":
        return _apply_pptx_operations(file_bytes, ops)
    raise ValueError(f"Unsupported office document type: {doc_type}")


def _extract_docx_state(file_bytes: bytes) -> dict[str, Any]:
    doc = Document(BytesIO(file_bytes))
    paragraphs: list[dict[str, Any]] = []
    headings: list[dict[str, Any]] = []
    for idx, paragraph in enumerate(doc.paragraphs):
        text = paragraph.text.strip()
        style_name = paragraph.style.name if paragraph.style else ""
        if not text:
            continue
        item = {"index": idx, "text": text, "style": style_name}
        paragraphs.append(item)
        if style_name.startswith("Heading"):
            headings.append(item)

    tables: list[dict[str, Any]] = []
    for table_idx, table in enumerate(doc.tables):
        tables.append(
            {
                "index": table_idx,
                "rows": [
                    [cell.text.strip() for cell in row.cells]
                    for row in table.rows
                ],
            }
        )

    return {"paragraphs": paragraphs[:250], "headings": headings[:50], "tables": tables[:20]}


def _apply_docx_operations(file_bytes: bytes, ops: list[dict[str, Any]]) -> bytes:
    doc = Document(BytesIO(file_bytes))
    for op in ops:
        op_type = op.get("type")
        if op_type == "add_heading":
            doc.add_heading(op.get("text", ""), level=max(0, min(int(op.get("level", 1)), 9)))
        elif op_type == "add_paragraph":
            paragraph = doc.add_paragraph(op.get("text", ""))
            style = op.get("style")
            if style:
                try:
                    paragraph.style = style
                except Exception:
                    pass
        elif op_type == "add_table":
            rows = max(int(op.get("rows", 1)), 1)
            cols = max(int(op.get("cols", 1)), 1)
            table = doc.add_table(rows=rows, cols=cols)
            data = op.get("data", [])
            if isinstance(data, list):
                for row_idx, row in enumerate(data[:rows]):
                    if not isinstance(row, list):
                        continue
                    for col_idx, value in enumerate(row[:cols]):
                        table.cell(row_idx, col_idx).text = str(value)
        elif op_type == "replace_text":
            find_text = str(op.get("find", ""))
            replace_text = str(op.get("replace", ""))
            if find_text:
                for paragraph in doc.paragraphs:
                    if find_text in paragraph.text:
                        paragraph.text = paragraph.text.replace(find_text, replace_text)
        elif op_type == "set_style":
            idx = int(op.get("paragraph_index", -1))
            if 0 <= idx < len(doc.paragraphs):
                try:
                    doc.paragraphs[idx].style = op.get("style", "Normal")
                except Exception:
                    pass
        elif op_type == "add_page_break":
            doc.add_paragraph().add_run().add_break(WD_BREAK.PAGE)
        elif op_type == "clear_after":
            idx = int(op.get("paragraph_index", -1))
            for paragraph in list(doc.paragraphs[idx + 1 :]):
                _delete_docx_paragraph(paragraph)
        elif op_type == "bold_text":
            idx = int(op.get("paragraph_index", -1))
            if 0 <= idx < len(doc.paragraphs):
                for run in doc.paragraphs[idx].runs:
                    run.bold = True
        elif op_type == "set_font":
            idx = int(op.get("paragraph_index", -1))
            if 0 <= idx < len(doc.paragraphs):
                font_name = op.get("name")
                font_size = op.get("size")
                for run in doc.paragraphs[idx].runs:
                    if font_name:
                        run.font.name = str(font_name)
                    if font_size:
                        run.font.size = _docx_pt(float(font_size))

    output = BytesIO()
    doc.save(output)
    return output.getvalue()


def _delete_docx_paragraph(paragraph) -> None:
    element = paragraph._element
    parent = element.getparent()
    if parent is not None:
        parent.remove(element)


def _docx_pt(value: float):
    from docx.shared import Pt

    return Pt(value)


def _extract_xlsx_state(file_bytes: bytes) -> dict[str, Any]:
    workbook = load_workbook(BytesIO(file_bytes))
    sheets: list[dict[str, Any]] = []
    for sheet in workbook.worksheets:
        rows: list[list[str]] = []
        for row in sheet.iter_rows(min_row=1, max_row=min(sheet.max_row, 100), values_only=True):
            rows.append(["" if cell is None else str(cell) for cell in row[:25]])
        sheets.append({"name": sheet.title, "rows": rows})

    named_ranges = [name.name for name in workbook.defined_names.definedName]
    return {"sheet_names": workbook.sheetnames, "named_ranges": named_ranges, "sheets": sheets}


def _apply_xlsx_operations(file_bytes: bytes, ops: list[dict[str, Any]]) -> bytes:
    workbook = load_workbook(BytesIO(file_bytes))
    for op in ops:
        op_type = op.get("type")
        if op_type == "add_sheet":
            name = str(op.get("name", "Sheet"))
            if name not in workbook.sheetnames:
                workbook.create_sheet(title=name)
        elif op_type == "rename_sheet":
            old = str(op.get("old", ""))
            new = str(op.get("new", ""))
            if old in workbook.sheetnames and new:
                workbook[old].title = new
        else:
            sheet_name = str(op.get("sheet", workbook.sheetnames[0]))
            sheet = workbook[sheet_name] if sheet_name in workbook.sheetnames else workbook.active
            row = int(op.get("row", 1))
            col = int(op.get("col", 1))
            if op_type == "set_cell":
                sheet.cell(row=row, column=col, value=op.get("value"))
            elif op_type == "set_formula":
                sheet.cell(row=row, column=col, value=op.get("formula"))
            elif op_type == "add_row":
                row_index = int(op.get("row_index", row))
                sheet.insert_rows(row_index)
                for idx, value in enumerate(op.get("data", []), start=1):
                    sheet.cell(row=row_index, column=idx, value=value)
            elif op_type == "bold_row":
                bold_font = Font(bold=True)
                for cell in sheet[int(op.get("row", row))]:
                    cell.font = bold_font
            elif op_type == "set_fill":
                color = str(op.get("hex_color", "D4A017")).replace("#", "")
                sheet.cell(row=row, column=col).fill = PatternFill(
                    fill_type="solid",
                    fgColor=color,
                )
            elif op_type == "auto_fit_columns":
                _auto_fit_sheet(sheet)

    output = BytesIO()
    workbook.save(output)
    return output.getvalue()


def _auto_fit_sheet(sheet) -> None:
    for idx in range(1, sheet.max_column + 1):
        max_length = 0
        for cell in sheet[get_column_letter(idx)]:
            max_length = max(max_length, len(str(cell.value or "")))
        sheet.column_dimensions[get_column_letter(idx)].width = min(max(max_length + 2, 10), 40)


def _extract_pptx_state(file_bytes: bytes) -> dict[str, Any]:
    presentation = Presentation(BytesIO(file_bytes))
    slides: list[dict[str, Any]] = []
    for idx, slide in enumerate(presentation.slides):
        title = slide.shapes.title.text.strip() if slide.shapes.title and slide.shapes.title.text else ""
        bullets: list[str] = []
        shape_types: list[str] = []
        for shape in slide.shapes:
            shape_types.append(str(shape.shape_type))
            if not hasattr(shape, "text_frame") or not shape.text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                text = paragraph.text.strip()
                if text and text != title:
                    bullets.append(text)
        slides.append(
            {
                "index": idx,
                "title": title,
                "bullets": bullets[:20],
                "shape_types": shape_types,
            }
        )
    return {"slide_count": len(slides), "slides": slides}


def _extract_pdf_state(file_bytes: bytes) -> dict[str, Any]:
    try:
        import fitz

        pdf = fitz.open(stream=file_bytes, filetype="pdf")
        pages: list[dict[str, Any]] = []
        for index, page in enumerate(pdf, start=1):
            text = page.get_text("text").strip()
            pages.append(
                {
                    "index": index - 1,
                    "page_number": index,
                    "text": text[:2000],
                }
            )
        pdf.close()
        return {
            "page_count": len(pages),
            "pages": pages[:20],
        }
    except ImportError:
        return _extract_pdf_state_fallback(file_bytes)


def _apply_pptx_operations(file_bytes: bytes, ops: list[dict[str, Any]]) -> bytes:
    presentation = Presentation(BytesIO(file_bytes))
    for op in ops:
        op_type = op.get("type")
        if op_type == "add_slide":
            layout_index = int(op.get("layout_index", 1))
            layout = presentation.slide_layouts[min(layout_index, len(presentation.slide_layouts) - 1)]
            slide = presentation.slides.add_slide(layout)
            if slide.shapes.title:
                slide.shapes.title.text = str(op.get("title", ""))
            body = op.get("body")
            if body and len(slide.placeholders) > 1:
                slide.placeholders[1].text = str(body)
        elif op_type == "edit_slide_title":
            slide = _get_slide(presentation, int(op.get("slide_index", -1)))
            if slide and slide.shapes.title:
                slide.shapes.title.text = str(op.get("text", ""))
        elif op_type == "edit_slide_body":
            slide = _get_slide(presentation, int(op.get("slide_index", -1)))
            if slide:
                placeholder = _body_placeholder(slide)
                if placeholder is not None:
                    placeholder.text = str(op.get("text", ""))
        elif op_type == "add_bullet":
            slide = _get_slide(presentation, int(op.get("slide_index", -1)))
            if slide:
                placeholder = _body_placeholder(slide)
                if placeholder is not None:
                    paragraph = placeholder.text_frame.add_paragraph()
                    paragraph.text = str(op.get("text", ""))
                    paragraph.level = int(op.get("level", 0))
        elif op_type == "add_image_placeholder":
            slide = _get_slide(presentation, int(op.get("slide_index", -1)))
            if slide:
                textbox = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(4), Inches(0.6))
                textbox.text_frame.text = "[Image placeholder]"
        elif op_type == "reorder_slides":
            new_order = op.get("new_order", [])
            if isinstance(new_order, list):
                _reorder_slides(presentation, [int(idx) for idx in new_order])
        elif op_type == "delete_slide":
            _delete_slide(presentation, int(op.get("slide_index", -1)))
        elif op_type == "set_slide_notes":
            slide = _get_slide(presentation, int(op.get("slide_index", -1)))
            if slide:
                slide.notes_slide.notes_text_frame.text = str(op.get("text", ""))

    output = BytesIO()
    presentation.save(output)
    return output.getvalue()


def _get_slide(presentation: Presentation, index: int):
    if 0 <= index < len(presentation.slides):
        return presentation.slides[index]
    return None


def _body_placeholder(slide):
    for shape in slide.placeholders:
        if shape.is_placeholder and shape.placeholder_format.type == PP_PLACEHOLDER.BODY:
            return shape
    if len(slide.placeholders) > 1:
        return slide.placeholders[1]
    return None


def _delete_slide(presentation: Presentation, index: int) -> None:
    if not 0 <= index < len(presentation.slides):
        return
    slide_id_list = presentation.slides._sldIdLst
    slides = list(slide_id_list)
    slide_id_list.remove(slides[index])


def _reorder_slides(presentation: Presentation, order: list[int]) -> None:
    slide_id_list = presentation.slides._sldIdLst
    slides = list(slide_id_list)
    if sorted(order) != list(range(len(slides))):
        return
    for slide in list(slide_id_list):
        slide_id_list.remove(slide)
    for idx in order:
        slide_id_list.append(slides[idx])


async def _summarize_changes(
    instruction: str,
    doc_type: str,
    ops: list[dict[str, Any]],
) -> str:
    op_text = json.dumps(ops, ensure_ascii=True)
    response = await chat_completion(
        messages=[
            {
                "role": "system",
                "content": "Summarize document edits for the user in 1-2 concise sentences.",
            },
            {
                "role": "user",
                "content": (
                    f"Document type: {doc_type}\n"
                    f"Instruction: {instruction}\n"
                    f"Applied operations JSON: {op_text}"
                ),
            },
        ],
        tools=None,
        model="gpt-4o",
    )
    content = getattr(response.choices[0].message, "content", "") or ""
    if content and "stub mode" not in content.lower():
        return content
    if not ops:
        return "Amin reviewed the document but did not apply any concrete changes."
    changed = ", ".join(op.get("type", "update") for op in ops[:4])
    return f"Amin updated the {doc_type.upper()} document by applying: {changed}."


def _fallback_document_answer(
    state: dict[str, Any],
    doc_type: str,
    question: str,
) -> str:
    if doc_type == "docx":
        preview = [item.get("text", "") for item in state.get("paragraphs", [])[:5]]
        return f"From the current DOCX document, the visible structure includes: {' | '.join(preview) or 'no readable paragraphs found'}."
    if doc_type == "xlsx":
        sheet_names = ", ".join(state.get("sheet_names", []))
        return f"The workbook contains these sheets: {sheet_names or 'none detected'}. Question asked: {question}"
    if doc_type == "pptx":
        slide_titles = ", ".join(
            slide.get("title") or f"Slide {idx + 1}"
            for idx, slide in enumerate(state.get("slides", [])[:6])
        )
        return f"The presentation currently includes these slides: {slide_titles or 'none detected'}. Question asked: {question}"

    preview = " | ".join(
        page.get("text", "")[:180]
        for page in state.get("pages", [])[:3]
        if page.get("text")
    )
    return f"The PDF includes {state.get('page_count', 0)} pages. Visible preview: {preview or 'no readable text extracted'}."


def _extract_pdf_state_fallback(file_bytes: bytes) -> dict[str, Any]:
    raw_text = file_bytes.decode("latin-1", errors="ignore")
    matches = re.findall(r"\(([^()]*)\)\s*Tj", raw_text)
    joined = "\n".join(match.strip() for match in matches if match.strip())[:4000]
    return {
        "page_count": 1,
        "pages": [
            {
                "index": 0,
                "page_number": 1,
                "text": joined,
            }
        ],
    }
